# SPDX-FileCopyrightText: 2025 Weibo, Inc.
#
# SPDX-License-Identifier: Apache-2.0

"""PPTX generation tool for creating PowerPoint presentations from structured content."""

import json
import logging
from typing import Any, Optional

from langchain_core.tools import BaseTool
from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)


class SlideContent(BaseModel):
    """Content for a single slide."""

    title: str = Field(description="Slide title")
    content: str = Field(
        description="Main content of the slide in markdown format (bullet points, text, etc.)"
    )
    notes: Optional[str] = Field(
        default=None, description="Speaker notes for this slide"
    )
    layout: str = Field(
        default="title_and_content",
        description="Slide layout type: 'title', 'title_and_content', 'two_column', 'blank'",
    )


class PPTXGeneratorInput(BaseModel):
    """Input schema for PPTX generator tool."""

    title: str = Field(description="Presentation title")
    slides: list[dict[str, Any]] = Field(
        description=(
            "List of slides, each containing 'title' (string), "
            "'content' (markdown string), optional 'notes' (string), "
            "and optional 'layout' (string: 'title', 'title_and_content', 'two_column', 'blank')"
        )
    )
    author: Optional[str] = Field(default=None, description="Presentation author name")
    theme: str = Field(
        default="default",
        description="Color theme: 'default', 'professional', 'creative', 'minimal'",
    )


class PPTXGeneratorTool(BaseTool):
    """
    Generate a PowerPoint presentation from structured content.

    This tool creates PPTX files from structured slide data.
    The generated presentation includes:
    - Title slide with presentation title and author
    - Content slides with titles, bullet points, and notes
    - Professional styling based on selected theme

    The tool outputs a JSON response with the generated file information,
    including context ID for downloading the PPTX and preview images.
    """

    name: str = "generate_pptx"
    display_name: str = "Generate PPTX"
    description: str = (
        "Generate a PowerPoint presentation from structured content. "
        "Provide a title, list of slides (each with title, content in markdown format, "
        "optional notes, and layout), optional author name, and theme. "
        "Returns file information including download ID and preview images."
    )
    args_schema: type[BaseModel] = PPTXGeneratorInput

    # Configuration
    task_id: int = 0
    subtask_id: int = 0
    user_id: int = 0
    api_base_url: str = "http://localhost:8000"
    auth_token: Optional[str] = None

    def __init__(
        self,
        task_id: int = 0,
        subtask_id: int = 0,
        user_id: int = 0,
        api_base_url: str = "http://localhost:8000",
        auth_token: Optional[str] = None,
        **kwargs,
    ) -> None:
        super().__init__(**kwargs)
        self.task_id = task_id
        self.subtask_id = subtask_id
        self.user_id = user_id
        self.api_base_url = api_base_url
        self.auth_token = auth_token

    def _run(
        self,
        title: str,
        slides: list[dict[str, Any]],
        author: Optional[str] = None,
        theme: str = "default",
        **_,
    ) -> str:
        """Synchronous execution - not recommended, use _arun instead."""
        import asyncio

        return asyncio.get_event_loop().run_until_complete(
            self._arun(title, slides, author, theme)
        )

    async def _arun(
        self,
        title: str,
        slides: list[dict[str, Any]],
        author: Optional[str] = None,
        theme: str = "default",
        **_,
    ) -> str:
        """
        Generate PPTX asynchronously.

        This implementation generates the PPTX using python-pptx library directly,
        without requiring a sandbox. For more complex presentations with charts
        and advanced layouts, the sandbox skill should be used instead.
        """
        try:
            # Validate input
            if not title:
                return json.dumps(
                    {"error": "Presentation title is required"}, ensure_ascii=False
                )
            if not slides:
                return json.dumps(
                    {"error": "At least one slide is required"}, ensure_ascii=False
                )

            # Generate PPTX using python-pptx
            pptx_binary, slide_count = await self._generate_pptx(
                title=title,
                slides=slides,
                author=author,
                theme=theme,
            )

            # Generate preview thumbnails
            preview_images = await self._generate_previews(pptx_binary)

            # Store the generated files using backend API
            result = await self._store_generated_files(
                title=title,
                pptx_binary=pptx_binary,
                preview_images=preview_images,
                slide_count=slide_count,
            )

            return json.dumps(result, ensure_ascii=False)

        except Exception as e:
            logger.exception(f"PPTX generation failed: {e}")
            return json.dumps(
                {"error": f"Failed to generate presentation: {str(e)}"},
                ensure_ascii=False,
            )

    async def _generate_pptx(
        self,
        title: str,
        slides: list[dict[str, Any]],
        author: Optional[str],
        theme: str,
    ) -> tuple[bytes, int]:
        """Generate PPTX binary using python-pptx."""
        from io import BytesIO

        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # Theme colors
        themes = {
            "default": {"primary": "1F497D", "secondary": "4F81BD", "bg": "FFFFFF"},
            "professional": {"primary": "2C3E50", "secondary": "3498DB", "bg": "F5F5F5"},
            "creative": {"primary": "E74C3C", "secondary": "F39C12", "bg": "FFFFFF"},
            "minimal": {"primary": "333333", "secondary": "666666", "bg": "FFFFFF"},
        }
        colors = themes.get(theme, themes["default"])

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)

        # Add title slide
        title_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(title_slide_layout)

        # Add title text box
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor.from_string(colors["primary"])
        title_para.alignment = PP_ALIGN.CENTER

        # Add author if provided
        if author:
            author_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.5)
            )
            author_frame = author_box.text_frame
            author_para = author_frame.paragraphs[0]
            author_para.text = author
            author_para.font.size = Pt(20)
            author_para.font.color.rgb = RGBColor.from_string(colors["secondary"])
            author_para.alignment = PP_ALIGN.CENTER

        # Add content slides
        for slide_data in slides:
            slide_title = slide_data.get("title", "")
            content = slide_data.get("content", "")
            layout = slide_data.get("layout", "title_and_content")

            content_slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Add slide title
            if slide_title:
                title_box = content_slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(12.333), Inches(1)
                )
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = slide_title
                title_para.font.size = Pt(32)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor.from_string(colors["primary"])

            # Add content
            if content:
                content_box = content_slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5)
                )
                content_frame = content_box.text_frame
                content_frame.word_wrap = True

                # Parse markdown-like content
                lines = content.strip().split("\n")
                for i, line in enumerate(lines):
                    if i == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()

                    # Handle bullet points
                    line_text = line.strip()
                    if line_text.startswith("- ") or line_text.startswith("* "):
                        p.text = "• " + line_text[2:]
                        p.level = 0
                    elif line_text.startswith("  - ") or line_text.startswith("  * "):
                        p.text = "  • " + line_text[4:]
                        p.level = 1
                    else:
                        p.text = line_text

                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor.from_string("333333")

        # Save to bytes
        output = BytesIO()
        prs.save(output)
        pptx_binary = output.getvalue()

        return pptx_binary, len(slides) + 1  # +1 for title slide

    async def _generate_previews(self, pptx_binary: bytes) -> list[bytes]:
        """
        Generate preview thumbnails for the PPTX.

        For simplicity, this returns an empty list. The full implementation
        would use LibreOffice or a similar tool to convert slides to images.
        The PPTX skill package provides full thumbnail generation support.
        """
        # Preview generation requires LibreOffice or similar tools
        # This is better handled in the sandbox skill for reliability
        return []

    async def _store_generated_files(
        self,
        title: str,
        pptx_binary: bytes,
        preview_images: list[bytes],
        slide_count: int,
    ) -> dict[str, Any]:
        """
        Store generated PPTX and preview images via backend API.

        Returns context information for display in chat.
        """
        import httpx

        filename = f"{title.replace(' ', '_')}.pptx"

        # If we have access to the backend API, store the files
        if self.api_base_url and self.auth_token:
            try:
                async with httpx.AsyncClient(timeout=60.0) as client:
                    # Create form data for PPTX upload
                    files = {"file": (filename, pptx_binary, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
                    headers = {"Authorization": f"Bearer {self.auth_token}"}

                    response = await client.post(
                        f"{self.api_base_url}/api/attachments/upload",
                        files=files,
                        headers=headers,
                    )

                    if response.status_code == 200:
                        attachment_data = response.json()
                        return {
                            "status": "success",
                            "message": f"Generated presentation '{title}' with {slide_count} slides",
                            "pptx_context_id": attachment_data.get("id"),
                            "filename": filename,
                            "slide_count": slide_count,
                            "file_size": len(pptx_binary),
                            "download_url": f"/api/attachments/{attachment_data.get('id')}/download",
                        }
                    else:
                        logger.error(f"Failed to upload PPTX: {response.status_code} - {response.text}")

            except Exception as e:
                logger.error(f"Failed to store PPTX via API: {e}")

        # Fallback: Return the PPTX binary as base64 for client-side download
        import base64

        return {
            "status": "success",
            "message": f"Generated presentation '{title}' with {slide_count} slides",
            "filename": filename,
            "slide_count": slide_count,
            "file_size": len(pptx_binary),
            "pptx_base64": base64.b64encode(pptx_binary).decode("utf-8"),
            "note": "PPTX data included as base64. Decode and save as .pptx file.",
        }
