# SPDX-FileCopyrightText: 2025 Weibo, Inc.
#
# SPDX-License-Identifier: Apache-2.0

"""
Document parser service for extracting text from various file formats.

Supports: PDF, Word (.doc, .docx), PowerPoint (.ppt, .pptx),
Excel (.xls, .xlsx, .csv), TXT, Markdown files, Images (.jpg, .jpeg, .png, .gif, .bmp, .webp),
and any text-based files detected via MIME type analysis.

Features smart truncation that preserves document structure:
- Excel/CSV: Header + sample rows + ellipsis + tail rows
- PDF: First pages + middle summary + last pages
- Word: Opening paragraphs + middle summary + closing paragraphs
- PowerPoint: First/last slides + middle summary
- Text/Markdown: Head content + middle summary + tail content

MIME-based text file detection:
- Uses python-magic to detect actual file content type
- Supports text/* MIME types and common application/* text formats
- Allows uploading code files (.py, .js, .java, etc.) without explicit extension whitelist
"""

import base64
import csv
import io
import logging
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Set, Tuple

import chardet
import magic

from app.core.config import settings
from app.services.attachment.smart_truncation import (
    SmartTruncationConfig,
    SmartTruncationInfo,
    SmartTruncationManager,
    TruncationType,
)

logger = logging.getLogger(__name__)


# MIME types that are considered text-based files
# These files can be parsed as plain text
TEXT_MIME_TYPES: Set[str] = {
    # text/* family
    "text/plain",
    "text/html",
    "text/css",
    "text/javascript",
    "text/xml",
    "text/csv",
    "text/markdown",
    "text/x-python",
    "text/x-java",
    "text/x-c",
    "text/x-c++",
    "text/x-ruby",
    "text/x-perl",
    "text/x-php",
    "text/x-shellscript",
    "text/x-script.python",
    "text/x-go",
    "text/x-rust",
    "text/x-swift",
    "text/x-kotlin",
    "text/x-scala",
    "text/x-typescript",
    "text/x-coffeescript",
    "text/x-lua",
    "text/x-r",
    "text/x-matlab",
    "text/x-sql",
    "text/x-yaml",
    "text/x-toml",
    "text/x-ini",
    "text/x-properties",
    "text/x-diff",
    "text/x-patch",
    "text/x-log",
    "text/x-makefile",
    "text/x-cmake",
    "text/x-dockerfile",
    "text/x-nginx-conf",
    "text/x-apache-conf",
    "text/x-systemd-unit",
    "text/x-tex",
    "text/x-latex",
    "text/x-bibtex",
    "text/x-rst",
    "text/x-asciidoc",
    "text/x-org",
    "text/troff",
    "text/rtf",
    "text/calendar",
    "text/vcard",
    # application/* text-based types
    "application/json",
    "application/xml",
    "application/javascript",
    "application/x-javascript",
    "application/ecmascript",
    "application/x-sh",
    "application/x-bash",
    "application/x-csh",
    "application/x-zsh",
    "application/x-python",
    "application/x-ruby",
    "application/x-perl",
    "application/x-php",
    "application/sql",
    "application/graphql",
    "application/toml",
    "application/x-yaml",
    "application/yaml",
    "application/x-httpd-php",
    "application/x-typescript",
    "application/typescript",
    "application/x-tex",
    "application/x-latex",
    "application/x-troff",
    "application/x-troff-man",
    "application/x-ndjson",
    "application/ld+json",
    "application/manifest+json",
    "application/schema+json",
    "application/vnd.api+json",
    "application/hal+json",
    "application/problem+json",
    "application/x-www-form-urlencoded",
    "application/xhtml+xml",
    "application/atom+xml",
    "application/rss+xml",
    "application/soap+xml",
    "application/mathml+xml",
    "application/xslt+xml",
    "application/x-subrip",
    "application/x-wine-extension-ini",
}


@dataclass
class TruncationInfo:
    """Information about content truncation."""

    is_truncated: bool = False
    original_length: Optional[int] = None
    truncated_length: Optional[int] = None

    # Smart truncation details
    truncation_type: str = "none"  # "none", "simple", "smart"
    original_structure: Dict[str, Any] = field(default_factory=dict)
    kept_structure: Dict[str, Any] = field(default_factory=dict)
    summary_message: str = ""

    @classmethod
    def from_smart_info(cls, smart_info: SmartTruncationInfo) -> "TruncationInfo":
        """Create TruncationInfo from SmartTruncationInfo."""
        return cls(
            is_truncated=smart_info.is_truncated,
            original_length=smart_info.original_length,
            truncated_length=smart_info.truncated_length,
            truncation_type=smart_info.truncation_type.value,
            original_structure=smart_info.original_structure,
            kept_structure=smart_info.kept_structure,
            summary_message=smart_info.summary_message,
        )


@dataclass
class ParseResult:
    """Result of document parsing."""

    text: str
    text_length: int
    image_base64: Optional[str] = None
    truncation_info: Optional[TruncationInfo] = None


class DocumentParseError(Exception):
    """Exception raised when document parsing fails."""

    # Error codes for i18n mapping
    UNSUPPORTED_TYPE = "unsupported_type"
    UNRECOGNIZED_TYPE = "unrecognized_type"
    FILE_TOO_LARGE = "file_too_large"
    PARSE_FAILED = "parse_failed"
    ENCRYPTED_PDF = "encrypted_pdf"
    LEGACY_DOC = "legacy_doc"
    LEGACY_PPT = "legacy_ppt"
    LEGACY_XLS = "legacy_xls"

    def __init__(self, message: str, error_code: Optional[str] = None):
        super().__init__(message)
        self.error_code = error_code or self.PARSE_FAILED


class DocumentParser:
    """
    Document parser for extracting text content from various file formats.

    Supported formats:
    - PDF (.pdf)
    - Word (.doc, .docx)
    - PowerPoint (.ppt, .pptx)
    - Excel (.xls, .xlsx, .csv)
    - Plain text (.txt)
    - Markdown (.md)
    - Images (.jpg, .jpeg, .png, .gif, .bmp, .webp)
    - Any text-based files detected via MIME type analysis

    Features smart truncation that preserves document structure instead of
    simple text cutting.

    For files with unknown extensions, the parser uses MIME type detection
    to identify text-based files (code files, config files, etc.) and
    processes them as plain text.
    """

    # Supported file extensions and their MIME types (known formats)
    SUPPORTED_EXTENSIONS = {
        ".pdf": "application/pdf",
        ".doc": "application/msword",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".ppt": "application/vnd.ms-powerpoint",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xls": "application/vnd.ms-excel",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".csv": "text/csv",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
        ".webp": "image/webp",
    }

    # Special format extensions that have dedicated parsers
    SPECIAL_FORMAT_EXTENSIONS = {
        ".pdf",
        ".doc",
        ".docx",
        ".ppt",
        ".pptx",
        ".xls",
        ".xlsx",
        ".csv",
        ".jpg",
        ".jpeg",
        ".png",
        ".gif",
        ".bmp",
        ".webp",
    }

    # Known text format extensions (no MIME detection needed)
    KNOWN_TEXT_EXTENSIONS = {".txt", ".md"}

    def __init__(self, truncation_config: Optional[SmartTruncationConfig] = None):
        """
        Initialize DocumentParser with optional truncation configuration.

        Args:
            truncation_config: Configuration for smart truncation.
                              If None, uses default configuration.
        """
        self.truncation_manager = SmartTruncationManager(truncation_config)

    @classmethod
    def get_max_file_size(cls) -> int:
        """Get maximum file size from configuration (in bytes)."""
        return settings.MAX_UPLOAD_FILE_SIZE_MB * 1024 * 1024

    @classmethod
    def get_max_text_length(cls) -> int:
        """Get maximum extracted text length from configuration."""
        return settings.MAX_EXTRACTED_TEXT_LENGTH

    @classmethod
    def is_supported_extension(cls, extension: str) -> bool:
        """Check if the file extension is supported."""
        return extension.lower() in cls.SUPPORTED_EXTENSIONS

    @classmethod
    def get_mime_type(cls, extension: str) -> str:
        """Get MIME type for a file extension."""
        return cls.SUPPORTED_EXTENSIONS.get(
            extension.lower(), "application/octet-stream"
        )

    @classmethod
    def validate_file_size(cls, size: int) -> bool:
        """Check if file size is within limits."""
        return size <= cls.get_max_file_size()

    @classmethod
    def validate_text_length(cls, text: str) -> bool:
        """Check if extracted text length is within limits."""
        return len(text) <= cls.get_max_text_length()

    @staticmethod
    def detect_mime_type(binary_data: bytes) -> str:
        """
        Detect the MIME type of file content using python-magic.

        Args:
            binary_data: File binary content

        Returns:
            Detected MIME type string (e.g., 'text/plain', 'application/json')
        """
        try:
            mime = magic.Magic(mime=True)
            mime_type = mime.from_buffer(binary_data)
            return mime_type
        except Exception as e:
            logger.warning(f"Failed to detect MIME type: {e}")
            return "application/octet-stream"

    @staticmethod
    def is_text_mime_type(mime_type: Optional[str]) -> bool:
        """
        Check if a MIME type represents a text-based file.

        Args:
            mime_type: MIME type string to check, or None

        Returns:
            True if the MIME type is text-based, False otherwise
        """
        if not mime_type:
            return False

        # Check if it starts with text/
        if mime_type.startswith("text/"):
            return True

        # Check if it's in our whitelist of text-based application/* types
        if mime_type in TEXT_MIME_TYPES:
            return True

        # Check for common text-based patterns
        # Many text formats use +json, +xml suffixes
        if mime_type.endswith("+json") or mime_type.endswith("+xml"):
            return True

        return False

    @classmethod
    def is_supported_extension(cls, extension: str) -> bool:
        """
        Check if the file extension is supported.

        For known extensions, returns True if in SUPPORTED_EXTENSIONS.
        For unknown extensions, returns True to allow MIME-based detection.
        """
        ext = extension.lower()
        # Known extensions are always supported
        if ext in cls.SUPPORTED_EXTENSIONS:
            return True
        # For unknown extensions, we allow them to proceed
        # The parse() method will use MIME detection to validate
        return True

    @classmethod
    def is_known_extension(cls, extension: str) -> bool:
        """Check if the extension is in the known formats list."""
        return extension.lower() in cls.SUPPORTED_EXTENSIONS

    def parse(
        self,
        binary_data: bytes,
        extension: str,
        use_smart_truncation: bool = True,
    ) -> ParseResult:
        """
        Parse document and extract text content with smart truncation.

        Args:
            binary_data: File binary data
            extension: File extension (e.g., '.pdf', '.docx')
            use_smart_truncation: Whether to use smart truncation (default: True).
                                 If False, falls back to simple text cutting.

        Returns:
            ParseResult with extracted text, length, optional image_base64,
            and truncation_info if content was truncated

        Raises:
            DocumentParseError: If parsing fails
        """
        extension = extension.lower()

        try:
            image_base64 = None
            truncation_info = None
            max_length = self.get_max_text_length()

            # Check if this is a known special format with dedicated parser
            if extension in self.SPECIAL_FORMAT_EXTENSIONS:
                return self._parse_special_format(
                    binary_data, extension, use_smart_truncation, max_length
                )

            # Check if this is a known text format
            if extension in self.KNOWN_TEXT_EXTENSIONS:
                return self._parse_text_format(
                    binary_data, use_smart_truncation, max_length
                )

            # For unknown extensions, use MIME detection
            mime_type = self.detect_mime_type(binary_data)
            logger.info(
                f"Detected MIME type '{mime_type}' for file with extension '{extension}'"
            )

            if self.is_text_mime_type(mime_type):
                # Parse as text file
                return self._parse_text_format(
                    binary_data, use_smart_truncation, max_length
                )
            else:
                # MIME type is not recognized as text
                raise DocumentParseError(
                    f"Unrecognized file type: {extension} (detected MIME: {mime_type})",
                    DocumentParseError.UNRECOGNIZED_TYPE,
                )

        except DocumentParseError:
            raise
        except Exception as e:
            logger.error(f"Error parsing document: {e}", exc_info=True)
            raise DocumentParseError(
                f"Failed to parse document: {str(e)}",
                DocumentParseError.PARSE_FAILED,
            ) from e

    def _parse_special_format(
        self,
        binary_data: bytes,
        extension: str,
        use_smart_truncation: bool,
        max_length: int,
    ) -> ParseResult:
        """
        Parse files with special formats (PDF, Word, Excel, etc.).

        These formats have dedicated parsers.
        """
        image_base64 = None
        truncation_info = None

        if use_smart_truncation:
            if extension == ".pdf":
                text, truncation_info = self._parse_pdf_smart(binary_data, max_length)
            elif extension in [".doc", ".docx"]:
                text, truncation_info = self._parse_word_smart(
                    binary_data, extension, max_length
                )
            elif extension in [".ppt", ".pptx"]:
                text, truncation_info = self._parse_powerpoint_smart(
                    binary_data, extension, max_length
                )
            elif extension in [".xls", ".xlsx"]:
                text, truncation_info = self._parse_excel_smart(
                    binary_data, extension, max_length
                )
            elif extension == ".csv":
                text, truncation_info = self._parse_csv_smart(binary_data, max_length)
            elif extension in [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"]:
                text, image_base64 = self._parse_image(binary_data, extension)
            else:
                raise DocumentParseError(
                    f"Unsupported file type: {extension}",
                    DocumentParseError.UNSUPPORTED_TYPE,
                )
        else:
            # Fallback to simple parsing without smart truncation
            if extension == ".pdf":
                text = self._parse_pdf(binary_data)
            elif extension in [".doc", ".docx"]:
                text = self._parse_word(binary_data, extension)
            elif extension in [".ppt", ".pptx"]:
                text = self._parse_powerpoint(binary_data, extension)
            elif extension in [".xls", ".xlsx"]:
                text = self._parse_excel(binary_data, extension)
            elif extension == ".csv":
                text = self._parse_csv(binary_data)
            elif extension in [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"]:
                text, image_base64 = self._parse_image(binary_data, extension)
            else:
                raise DocumentParseError(
                    f"Unsupported file type: {extension}",
                    DocumentParseError.UNSUPPORTED_TYPE,
                )

            # Simple truncation for non-smart mode
            if len(text) > max_length:
                original_length = len(text)
                text = text[:max_length]
                truncation_info = TruncationInfo(
                    is_truncated=True,
                    original_length=original_length,
                    truncated_length=max_length,
                    truncation_type="simple",
                )
                logger.info(
                    f"Text truncated from {original_length} to {max_length} characters"
                )

        return ParseResult(
            text=text,
            text_length=len(text),
            image_base64=image_base64,
            truncation_info=truncation_info,
        )

    def _parse_text_format(
        self,
        binary_data: bytes,
        use_smart_truncation: bool,
        max_length: int,
    ) -> ParseResult:
        """
        Parse text-based files (txt, md, code files, config files, etc.).
        """
        truncation_info = None

        if use_smart_truncation:
            text, truncation_info = self._parse_text_smart(binary_data, max_length)
        else:
            text = self._parse_text(binary_data)
            # Simple truncation for non-smart mode
            if len(text) > max_length:
                original_length = len(text)
                text = text[:max_length]
                truncation_info = TruncationInfo(
                    is_truncated=True,
                    original_length=original_length,
                    truncated_length=max_length,
                    truncation_type="simple",
                )
                logger.info(
                    f"Text truncated from {original_length} to {max_length} characters"
                )

        return ParseResult(
            text=text,
            text_length=len(text),
            truncation_info=truncation_info,
        )

    def _parse_pdf(self, binary_data: bytes) -> str:
        """Parse PDF file and extract text."""
        try:
            from PyPDF2 import PdfReader

            pdf_file = io.BytesIO(binary_data)
            reader = PdfReader(pdf_file)

            # Check if PDF is encrypted
            if reader.is_encrypted:
                raise DocumentParseError(
                    "Cannot parse encrypted PDF file",
                    DocumentParseError.ENCRYPTED_PDF,
                )

            text_parts = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

            return "\n\n".join(text_parts)

        except DocumentParseError:
            raise
        except Exception as e:
            logger.error(f"Error parsing PDF: {e}", exc_info=True)
            raise DocumentParseError(
                f"Failed to parse PDF: {str(e)}",
                DocumentParseError.PARSE_FAILED,
            ) from e

    def _parse_word(self, binary_data: bytes, extension: str) -> str:
        """Parse Word document and extract text."""
        try:
            if extension == ".doc":
                # For .doc files, we need to handle them differently
                # python-docx only supports .docx format
                raise DocumentParseError(
                    "Legacy .doc format is not fully supported. Please convert to .docx",
                    DocumentParseError.LEGACY_DOC,
                )

            from docx import Document

            doc_file = io.BytesIO(binary_data)
            doc = Document(doc_file)

            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))

            return "\n\n".join(text_parts)

        except DocumentParseError:
            raise
        except Exception as e:
            logger.error(f"Error parsing Word document: {e}", exc_info=True)
            raise DocumentParseError(
                f"Failed to parse Word document: {str(e)}",
                DocumentParseError.PARSE_FAILED,
            ) from e

    def _parse_powerpoint(self, binary_data: bytes, extension: str) -> str:
        """Parse PowerPoint file and extract text (text only, no images)."""
        try:
            if extension == ".ppt":
                raise DocumentParseError(
                    "Legacy .ppt format is not fully supported. Please convert to .pptx",
                    DocumentParseError.LEGACY_PPT,
                )

            from pptx import Presentation

            ppt_file = io.BytesIO(binary_data)
            prs = Presentation(ppt_file)

            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))

                if len(slide_text) > 1:  # More than just the slide header
                    text_parts.append("\n".join(slide_text))

            return "\n\n".join(text_parts)

        except DocumentParseError:
            raise
        except Exception as e:
            logger.error(f"Error parsing PowerPoint: {e}", exc_info=True)
            raise DocumentParseError(
                f"Failed to parse PowerPoint: {str(e)}",
                DocumentParseError.PARSE_FAILED,
            ) from e

    def _parse_excel(self, binary_data: bytes, extension: str) -> str:
        """Parse Excel file and extract text."""
        try:
            if extension == ".xls":
                raise DocumentParseError(
                    "Legacy .xls format is not fully supported. Please convert to .xlsx",
                    DocumentParseError.LEGACY_XLS,
                )

            from openpyxl import load_workbook

            excel_file = io.BytesIO(binary_data)
            wb = load_workbook(excel_file, read_only=True, data_only=True)

            text_parts = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"--- Sheet: {sheet_name} ---"]

                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))
                    if row_values:
                        sheet_text.append(" | ".join(row_values))

                if len(sheet_text) > 1:
                    text_parts.append("\n".join(sheet_text))

            wb.close()
            return "\n\n".join(text_parts)

        except DocumentParseError:
            raise
        except Exception as e:
            logger.error(f"Error parsing Excel: {e}", exc_info=True)
            raise DocumentParseError(
                f"Failed to parse Excel: {str(e)}",
                DocumentParseError.PARSE_FAILED,
            ) from e

    def _parse_csv(self, binary_data: bytes) -> str:
        """Parse CSV file and extract text."""
        try:
            # Detect encoding
            detected = chardet.detect(binary_data)
            encoding = detected.get("encoding", "utf-8") or "utf-8"

            text = binary_data.decode(encoding)
            csv_file = io.StringIO(text)

            reader = csv.reader(csv_file)
            rows = []
            for row in reader:
                if any(cell.strip() for cell in row):
                    rows.append(" | ".join(cell.strip() for cell in row))

            return "\n".join(rows)

        except Exception as e:
            logger.error(f"Error parsing CSV: {e}", exc_info=True)
            raise DocumentParseError(
                f"Failed to parse CSV: {str(e)}",
                DocumentParseError.PARSE_FAILED,
            ) from e

    def _parse_text(self, binary_data: bytes) -> str:
        """Parse plain text or markdown file."""
        # Try common encodings in order of preference
        encodings_to_try = ["utf-8", "utf-8-sig", "gbk", "gb2312", "gb18030", "latin-1"]

        # First, try to detect encoding using chardet
        try:
            detected = chardet.detect(binary_data)
            detected_encoding = detected.get("encoding")
            if detected_encoding and detected_encoding.lower() not in [
                "ascii",
                "charmap",
            ]:
                # Insert detected encoding at the beginning if it's valid
                encodings_to_try.insert(0, detected_encoding)
        except Exception:
            pass  # Ignore chardet errors

        # Try each encoding
        last_error = None
        for encoding in encodings_to_try:
            try:
                return binary_data.decode(encoding)
            except (UnicodeDecodeError, LookupError) as e:
                last_error = e
                continue

        # If all encodings fail, try with error handling
        try:
            return binary_data.decode("utf-8", errors="replace")
        except Exception as e:
            logger.error(f"Error parsing text file: {e}", exc_info=True)
            raise DocumentParseError(
                f"Failed to parse text file: {str(last_error or e)}",
                DocumentParseError.PARSE_FAILED,
            ) from e

    def _parse_image(self, binary_data: bytes, extension: str) -> Tuple[str, str]:
        """
        Parse image file and return metadata information with base64 encoding.

        Returns:
            Tuple of (metadata_text, base64_encoded_image)
        """
        try:
            from PIL import Image

            image_file = io.BytesIO(binary_data)
            img = Image.open(image_file)

            # Extract image metadata
            width, height = img.size
            mode = img.mode
            format_name = img.format or extension[1:].upper()

            # Build description
            text = f"[图片文件]\n"
            text += f"格式: {format_name}\n"
            text += f"尺寸: {width} x {height} 像素\n"
            text += f"颜色模式: {mode}\n"
            text += f"文件大小: {len(binary_data)} 字节"

            # Try to extract EXIF data if available
            if hasattr(img, "_getexif") and img._getexif():
                text += "\n\n[EXIF 信息]"
                exif_data = img._getexif()
                if exif_data:
                    # Just mention EXIF is available, don't extract all
                    text += "\n包含 EXIF 元数据"

            # Encode image to base64 for vision models
            image_base64 = base64.b64encode(binary_data).decode("utf-8")

            return text, image_base64

        except Exception as e:
            logger.error(f"Error parsing image: {e}", exc_info=True)
            raise DocumentParseError(
                f"Failed to parse image: {str(e)}",
                DocumentParseError.PARSE_FAILED,
            ) from e

    # ==================== Smart Truncation Methods ====================

    def _parse_pdf_smart(
        self, binary_data: bytes, max_length: int
    ) -> Tuple[str, Optional[TruncationInfo]]:
        """
        Parse PDF file with smart page-based truncation.

        Keeps first N pages + last M pages, omits middle pages.
        """
        try:
            from PyPDF2 import PdfReader

            pdf_file = io.BytesIO(binary_data)
            reader = PdfReader(pdf_file)

            if reader.is_encrypted:
                raise DocumentParseError(
                    "Cannot parse encrypted PDF file",
                    DocumentParseError.ENCRYPTED_PDF,
                )

            # Extract text per page
            pages_text = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    pages_text.append(page_text)

            # Apply smart truncation
            text, smart_info = self.truncation_manager.truncate_pdf(
                pages_text, max_length
            )

            truncation_info = None
            if smart_info.is_truncated:
                truncation_info = TruncationInfo.from_smart_info(smart_info)

            return text, truncation_info

        except DocumentParseError:
            raise
        except Exception as e:
            logger.error(f"Error parsing PDF with smart truncation: {e}", exc_info=True)
            raise DocumentParseError(
                f"Failed to parse PDF: {str(e)}",
                DocumentParseError.PARSE_FAILED,
            ) from e

    def _parse_word_smart(
        self, binary_data: bytes, extension: str, max_length: int
    ) -> Tuple[str, Optional[TruncationInfo]]:
        """
        Parse Word document with smart paragraph-based truncation.

        Keeps first N paragraphs + last M paragraphs, omits middle.
        """
        try:
            if extension == ".doc":
                raise DocumentParseError(
                    "Legacy .doc format is not fully supported. Please convert to .docx",
                    DocumentParseError.LEGACY_DOC,
                )

            from docx import Document

            doc_file = io.BytesIO(binary_data)
            doc = Document(doc_file)

            # Extract paragraphs
            paragraphs = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text)

            # Also extract text from tables as separate paragraphs
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        paragraphs.append(" | ".join(row_text))

            # Apply smart truncation
            text, smart_info = self.truncation_manager.truncate_word(
                paragraphs, max_length
            )

            truncation_info = None
            if smart_info.is_truncated:
                truncation_info = TruncationInfo.from_smart_info(smart_info)

            return text, truncation_info

        except DocumentParseError:
            raise
        except Exception as e:
            logger.error(
                f"Error parsing Word with smart truncation: {e}", exc_info=True
            )
            raise DocumentParseError(
                f"Failed to parse Word document: {str(e)}",
                DocumentParseError.PARSE_FAILED,
            ) from e

    def _parse_powerpoint_smart(
        self, binary_data: bytes, extension: str, max_length: int
    ) -> Tuple[str, Optional[TruncationInfo]]:
        """
        Parse PowerPoint with smart slide-based truncation.

        Keeps first N slides + last M slides, omits middle.
        """
        try:
            if extension == ".ppt":
                raise DocumentParseError(
                    "Legacy .ppt format is not fully supported. Please convert to .pptx",
                    DocumentParseError.LEGACY_PPT,
                )

            from pptx import Presentation

            ppt_file = io.BytesIO(binary_data)
            prs = Presentation(ppt_file)

            # Extract text per slide
            slides_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"--- Slide {slide_num} ---"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)

                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_content.append(" | ".join(row_text))

                if len(slide_content) > 1:
                    slides_text.append("\n".join(slide_content))

            # Apply smart truncation
            text, smart_info = self.truncation_manager.truncate_powerpoint(
                slides_text, max_length
            )

            truncation_info = None
            if smart_info.is_truncated:
                truncation_info = TruncationInfo.from_smart_info(smart_info)

            return text, truncation_info

        except DocumentParseError:
            raise
        except Exception as e:
            logger.error(
                f"Error parsing PowerPoint with smart truncation: {e}", exc_info=True
            )
            raise DocumentParseError(
                f"Failed to parse PowerPoint: {str(e)}",
                DocumentParseError.PARSE_FAILED,
            ) from e

    def _parse_excel_smart(
        self, binary_data: bytes, extension: str, max_length: int
    ) -> Tuple[str, Optional[TruncationInfo]]:
        """
        Parse Excel file with smart row-based truncation.

        Keeps header rows + sample rows + tail rows, omits middle.
        """
        try:
            if extension == ".xls":
                raise DocumentParseError(
                    "Legacy .xls format is not fully supported. Please convert to .xlsx",
                    DocumentParseError.LEGACY_XLS,
                )

            from openpyxl import load_workbook

            excel_file = io.BytesIO(binary_data)
            wb = load_workbook(excel_file, read_only=True, data_only=True)

            # Extract data per sheet
            sheets_data = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []

                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        row_values.append(cell.value)
                    # Keep row even if all values are None (to preserve structure)
                    if any(v is not None for v in row_values):
                        rows.append(row_values)

                if rows:
                    sheets_data.append(
                        {
                            "name": sheet_name,
                            "rows": rows,
                        }
                    )

            wb.close()

            # Apply smart truncation
            text, smart_info = self.truncation_manager.truncate_excel(
                sheets_data, max_length
            )

            truncation_info = None
            if smart_info.is_truncated:
                truncation_info = TruncationInfo.from_smart_info(smart_info)

            return text, truncation_info

        except DocumentParseError:
            raise
        except Exception as e:
            logger.error(
                f"Error parsing Excel with smart truncation: {e}", exc_info=True
            )
            raise DocumentParseError(
                f"Failed to parse Excel: {str(e)}",
                DocumentParseError.PARSE_FAILED,
            ) from e

    def _parse_csv_smart(
        self, binary_data: bytes, max_length: int
    ) -> Tuple[str, Optional[TruncationInfo]]:
        """
        Parse CSV file with smart row-based truncation.

        Keeps header row + sample rows + tail rows, omits middle.
        """
        try:
            # Detect encoding
            detected = chardet.detect(binary_data)
            encoding = detected.get("encoding", "utf-8") or "utf-8"

            text = binary_data.decode(encoding)
            csv_file = io.StringIO(text)

            reader = csv.reader(csv_file)
            rows = []
            for row in reader:
                if any(cell.strip() for cell in row):
                    rows.append([cell.strip() for cell in row])

            # Apply smart truncation
            text, smart_info = self.truncation_manager.truncate_csv(rows, max_length)

            truncation_info = None
            if smart_info.is_truncated:
                truncation_info = TruncationInfo.from_smart_info(smart_info)

            return text, truncation_info

        except Exception as e:
            logger.error(f"Error parsing CSV with smart truncation: {e}", exc_info=True)
            raise DocumentParseError(
                f"Failed to parse CSV: {str(e)}",
                DocumentParseError.PARSE_FAILED,
            ) from e

    def _parse_text_smart(
        self, binary_data: bytes, max_length: int
    ) -> Tuple[str, Optional[TruncationInfo]]:
        """
        Parse text/markdown file with smart line-based truncation.

        Keeps first N lines + last M lines, omits middle.
        """
        # First decode the text
        text = self._parse_text(binary_data)

        # Apply smart truncation
        truncated_text, smart_info = self.truncation_manager.truncate_text(
            text, max_length
        )

        truncation_info = None
        if smart_info.is_truncated:
            truncation_info = TruncationInfo.from_smart_info(smart_info)

        return truncated_text, truncation_info


# Global parser instance
document_parser = DocumentParser()
