# SPDX-FileCopyrightText: 2025 WeCode, Inc.
#
# SPDX-License-Identifier: Apache-2.0

"""PPTX slide creation tool using E2B sandbox.

This tool allows Chat Shell agents to create PowerPoint presentations
by executing python-pptx code in an isolated E2B sandbox environment.

E2B Sandbox Flow:
1. Create sandbox using E2B SDK
2. Install python-pptx if not already installed
3. Execute PPTX generation code
4. Download the generated file from sandbox
5. Store as attachment in subtask_contexts table
6. Return download URL to user
"""

import asyncio
import json
import logging
import os
from typing import Any, Optional

from langchain_core.callbacks import CallbackManagerForToolRun
from langchain_core.tools import BaseTool
from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)

# Default configuration
DEFAULT_TIMEOUT = 120  # seconds
DEFAULT_MAX_RETRIES = 3
DEFAULT_OUTPUT_PATH = "/tmp/presentation.pptx"


class CreatePPTXInput(BaseModel):
    """Input schema for create_pptx tool."""

    content: str = Field(
        ...,
        description=(
            "The PPT content description or outline. Can be structured text with "
            "titles, bullet points, and sections. Use clear formatting like:\n"
            "Title Slide: [Title]\n"
            "Subtitle: [Subtitle]\n"
            "Slide 1: [Title]\n"
            "- Bullet point 1\n"
            "- Bullet point 2"
        ),
    )
    title: Optional[str] = Field(
        default=None,
        description="The main title of the presentation (used for title slide and filename).",
    )
    template: Optional[str] = Field(
        default="default",
        description=(
            "Template style to use: 'default', 'professional', 'minimal', "
            "or custom Python code for advanced templates."
        ),
    )
    filename: Optional[str] = Field(
        default="presentation.pptx",
        description="Output filename for the generated presentation.",
    )


class CreatePPTXTool(BaseTool):
    """Tool for creating PowerPoint presentations using E2B sandbox.

    This tool generates PPTX files by:
    1. Parsing content outline into slide structure
    2. Generating python-pptx code
    3. Executing code in E2B sandbox
    4. Downloading and storing the generated file
    5. Returning download URL

    Includes automatic retry with error correction for robustness.
    """

    name: str = "create_pptx"
    display_name: str = "Generate PPT"
    description: str = """Create a PowerPoint presentation (PPTX file) from content outline.

Use this tool when you need to generate professional slide presentations.

Parameters:
- content (required): PPT content outline with slides and bullet points
- title (optional): Main presentation title
- template (optional): Style template ("default", "professional", "minimal")
- filename (optional): Output filename (default: "presentation.pptx")

Returns:
- On success: attachment_id and download_url for the generated file
- On failure: Error message with suggestions

The generated file is automatically stored and can be downloaded via the provided URL."""

    args_schema: type[BaseModel] = CreatePPTXInput

    # Injected dependencies - set when creating the tool instance
    task_id: int = 0
    subtask_id: int = 0
    ws_emitter: Any = None
    user_id: int = 0
    user_name: str = ""

    # Configuration
    timeout: int = DEFAULT_TIMEOUT
    max_retries: int = DEFAULT_MAX_RETRIES

    class Config:
        arbitrary_types_allowed = True

    def _run(
        self,
        content: str,
        title: Optional[str] = None,
        template: Optional[str] = "default",
        filename: Optional[str] = "presentation.pptx",
        run_manager: CallbackManagerForToolRun | None = None,
    ) -> str:
        """Synchronous run - not implemented."""
        raise NotImplementedError("CreatePPTXTool only supports async execution")

    async def _arun(
        self,
        content: str,
        title: Optional[str] = None,
        template: Optional[str] = "default",
        filename: Optional[str] = "presentation.pptx",
        run_manager: CallbackManagerForToolRun | None = None,
    ) -> str:
        """Execute PPTX generation using E2B sandbox.

        Args:
            content: PPT content outline
            title: Optional presentation title
            template: Template style
            filename: Output filename
            run_manager: Callback manager

        Returns:
            JSON string with result (success with download URL or error)
        """
        logger.info(
            f"[PPTXTool] Starting PPTX generation: task_id={self.task_id}, "
            f"subtask_id={self.subtask_id}, content_length={len(content)}, "
            f"title={title}, template={template}"
        )

        # Ensure filename has .pptx extension
        if filename and not filename.endswith(".pptx"):
            filename = f"{filename}.pptx"
        effective_filename = filename or "presentation.pptx"

        # Emit status update via WebSocket
        await self._emit_tool_status("running", "Starting PPT generation...")

        try:
            # Run with retry mechanism
            result = await self._run_with_retry(
                content=content,
                title=title,
                template=template,
                filename=effective_filename,
            )

            # Check result and emit appropriate status
            if result.get("success"):
                await self._emit_tool_status(
                    "completed",
                    f"PPT created successfully: {effective_filename}",
                    result,
                )
            else:
                await self._emit_tool_status(
                    "failed",
                    result.get("error", "PPT generation failed"),
                    result,
                )

            return json.dumps(result, ensure_ascii=False, indent=2)

        except Exception as e:
            logger.error(
                f"[PPTXTool] Unexpected error: {e}",
                exc_info=True,
            )
            error_result = {
                "success": False,
                "error": f"Unexpected error: {e!s}",
                "suggestion": (
                    "An unexpected error occurred during PPT generation. "
                    "Please try again with a simpler content structure."
                ),
            }
            await self._emit_tool_status("failed", str(e))
            return json.dumps(error_result, ensure_ascii=False, indent=2)

    async def _run_with_retry(
        self,
        content: str,
        title: Optional[str],
        template: str,
        filename: str,
    ) -> dict:
        """Execute PPTX generation with retry mechanism.

        Args:
            content: PPT content outline
            title: Presentation title
            template: Template style
            filename: Output filename

        Returns:
            Result dict with success status and attachment info or error
        """
        last_error = None
        last_code = None

        for attempt in range(self.max_retries):
            logger.info(
                f"[PPTXTool] Generation attempt {attempt + 1}/{self.max_retries}"
            )

            try:
                # Generate python-pptx code
                pptx_code = self._generate_pptx_code(
                    content=content,
                    title=title,
                    template=template,
                    output_path=DEFAULT_OUTPUT_PATH,
                    previous_error=last_error,
                    previous_code=last_code,
                )
                last_code = pptx_code

                # Execute in E2B sandbox
                file_bytes = await self._execute_in_sandbox(pptx_code)

                if file_bytes:
                    # Store as attachment
                    attachment_result = await self._store_attachment(
                        file_bytes=file_bytes,
                        filename=filename,
                    )

                    if attachment_result.get("success"):
                        # Count slides for informational purposes
                        slide_count = self._estimate_slide_count(content)
                        return {
                            "success": True,
                            "attachment_id": attachment_result.get("attachment_id"),
                            "filename": filename,
                            "download_url": attachment_result.get("download_url"),
                            "slide_count": slide_count,
                            "message": (
                                f"PowerPoint presentation '{filename}' created successfully "
                                f"with approximately {slide_count} slides. "
                                f"The file has been saved and can be downloaded from the provided URL."
                            ),
                        }
                    else:
                        last_error = attachment_result.get(
                            "error", "Failed to store attachment"
                        )
                else:
                    last_error = "No file generated from sandbox execution"

            except Exception as e:
                last_error = str(e)
                logger.warning(f"[PPTXTool] Attempt {attempt + 1} failed: {last_error}")

        # All retries exhausted
        logger.error(
            f"[PPTXTool] All {self.max_retries} attempts failed. Last error: {last_error}"
        )
        return {
            "success": False,
            "error": f"Failed to generate presentation after {self.max_retries} attempts",
            "last_error": last_error,
            "final_instruction": (
                "CRITICAL: All automatic generation attempts have failed. "
                "DO NOT attempt to generate this presentation again. "
                "Instead, inform the user about the failure and suggest: "
                "1) Simplifying the content structure, "
                "2) Reducing the number of slides, "
                "3) Using clearer bullet point formatting."
            ),
        }

    def _generate_pptx_code(
        self,
        content: str,
        title: Optional[str],
        template: str,
        output_path: str,
        previous_error: Optional[str] = None,
        previous_code: Optional[str] = None,
    ) -> str:
        """Generate python-pptx code for creating the presentation.

        Args:
            content: PPT content outline
            title: Presentation title
            template: Template style
            output_path: Output file path in sandbox
            previous_error: Error from previous attempt (for auto-correction)
            previous_code: Code from previous attempt (for reference)

        Returns:
            Python code string to execute in sandbox
        """
        # Parse content into slides structure
        slides = self._parse_content_to_slides(content, title)

        # Get template colors based on style
        colors = self._get_template_colors(template)

        # Build the Python code
        code_lines = [
            "# Auto-generated PPTX creation code",
            "from pptx import Presentation",
            "from pptx.util import Inches, Pt",
            "from pptx.dml.color import RgbColor",
            "from pptx.enum.text import PP_ALIGN, MSO_ANCHOR",
            "",
            "# Create presentation",
            "prs = Presentation()",
            "prs.slide_width = Inches(13.333)",
            "prs.slide_height = Inches(7.5)",
            "",
            f"# Template colors: {template}",
            f"TITLE_COLOR = RgbColor({colors['title'][0]}, {colors['title'][1]}, {colors['title'][2]})",
            f"TEXT_COLOR = RgbColor({colors['text'][0]}, {colors['text'][1]}, {colors['text'][2]})",
            f"ACCENT_COLOR = RgbColor({colors['accent'][0]}, {colors['accent'][1]}, {colors['accent'][2]})",
            "",
        ]

        # Add helper function for text formatting
        code_lines.extend(
            [
                "def set_text_style(run, font_size, color, bold=False):",
                "    run.font.size = Pt(font_size)",
                "    run.font.color.rgb = color",
                "    run.font.bold = bold",
                "",
            ]
        )

        # Generate code for each slide
        for i, slide in enumerate(slides):
            slide_type = slide.get("type", "content")
            slide_title = slide.get("title", f"Slide {i + 1}")
            slide_content = slide.get("content", [])

            if slide_type == "title":
                # Title slide
                code_lines.extend(
                    [
                        f"# Slide {i + 1}: Title Slide",
                        "slide_layout = prs.slide_layouts[6]  # Blank layout",
                        "slide = prs.slides.add_slide(slide_layout)",
                        "",
                        "# Add title",
                        "title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))",
                        "title_frame = title_box.text_frame",
                        "title_para = title_frame.paragraphs[0]",
                        "title_para.alignment = PP_ALIGN.CENTER",
                        "title_run = title_para.add_run()",
                        f'title_run.text = """{self._escape_string(slide_title)}"""',
                        "set_text_style(title_run, 44, TITLE_COLOR, bold=True)",
                        "",
                    ]
                )

                # Add subtitle if present
                if slide_content:
                    subtitle = (
                        slide_content[0]
                        if isinstance(slide_content[0], str)
                        else str(slide_content[0])
                    )
                    code_lines.extend(
                        [
                            "# Add subtitle",
                            "subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))",
                            "subtitle_frame = subtitle_box.text_frame",
                            "subtitle_para = subtitle_frame.paragraphs[0]",
                            "subtitle_para.alignment = PP_ALIGN.CENTER",
                            "subtitle_run = subtitle_para.add_run()",
                            f'subtitle_run.text = """{self._escape_string(subtitle)}"""',
                            "set_text_style(subtitle_run, 24, TEXT_COLOR)",
                            "",
                        ]
                    )
            else:
                # Content slide
                code_lines.extend(
                    [
                        f"# Slide {i + 1}: Content Slide",
                        "slide_layout = prs.slide_layouts[6]  # Blank layout",
                        "slide = prs.slides.add_slide(slide_layout)",
                        "",
                        "# Add slide title",
                        "title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))",
                        "title_frame = title_box.text_frame",
                        "title_para = title_frame.paragraphs[0]",
                        "title_run = title_para.add_run()",
                        f'title_run.text = """{self._escape_string(slide_title)}"""',
                        "set_text_style(title_run, 32, TITLE_COLOR, bold=True)",
                        "",
                    ]
                )

                # Add content bullets
                if slide_content:
                    code_lines.extend(
                        [
                            "# Add content",
                            "content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.7), Inches(11.833), Inches(5.3))",
                            "content_frame = content_box.text_frame",
                            "content_frame.word_wrap = True",
                            "",
                        ]
                    )

                    for j, bullet in enumerate(slide_content):
                        bullet_text = bullet if isinstance(bullet, str) else str(bullet)
                        # Check for sub-items (indented with spaces or tabs)
                        is_sub_item = bullet_text.startswith(
                            "  "
                        ) or bullet_text.startswith("\t")
                        bullet_text = bullet_text.strip()

                        if j == 0:
                            code_lines.extend(
                                [
                                    "para = content_frame.paragraphs[0]",
                                ]
                            )
                        else:
                            code_lines.extend(
                                [
                                    "para = content_frame.add_paragraph()",
                                ]
                            )

                        font_size = "18" if is_sub_item else "20"
                        code_lines.extend(
                            [
                                f"para.level = {1 if is_sub_item else 0}",
                                "para.space_before = Pt(8)",
                                "para.space_after = Pt(4)",
                                "run = para.add_run()",
                                f'run.text = "• " + """{self._escape_string(bullet_text)}"""',
                                f"set_text_style(run, {font_size}, TEXT_COLOR)",
                                "",
                            ]
                        )

                code_lines.append("")

        # Save the presentation
        code_lines.extend(
            [
                "# Save presentation",
                f'prs.save("{output_path}")',
                f'print("SUCCESS: Presentation saved to {output_path}")',
            ]
        )

        return "\n".join(code_lines)

    def _parse_content_to_slides(
        self, content: str, title: Optional[str]
    ) -> list[dict]:
        """Parse content string into slide structure.

        Args:
            content: Content outline string
            title: Optional main title

        Returns:
            List of slide dicts with type, title, and content
        """
        slides = []
        current_slide = None
        has_title_slide = False

        lines = content.strip().split("\n")

        for line in lines:
            line_stripped = line.strip()
            if not line_stripped:
                continue

            # Check for title slide
            if line_stripped.lower().startswith("title slide:"):
                has_title_slide = True
                if current_slide:
                    slides.append(current_slide)
                slide_title = line_stripped[12:].strip()
                current_slide = {
                    "type": "title",
                    "title": slide_title or title or "Presentation",
                    "content": [],
                }
            # Check for subtitle (goes with title slide)
            elif line_stripped.lower().startswith("subtitle:"):
                if current_slide and current_slide.get("type") == "title":
                    subtitle = line_stripped[9:].strip()
                    current_slide["content"].append(subtitle)
            # Check for slide header
            elif line_stripped.lower().startswith("slide"):
                if current_slide:
                    slides.append(current_slide)
                # Extract slide title (e.g., "Slide 1: Overview" -> "Overview")
                parts = line_stripped.split(":", 1)
                slide_title = (
                    parts[1].strip() if len(parts) > 1 else f"Slide {len(slides) + 1}"
                )
                current_slide = {
                    "type": "content",
                    "title": slide_title,
                    "content": [],
                }
            # Check for bullet points
            elif line_stripped.startswith("-") or line_stripped.startswith("•"):
                bullet_text = line_stripped[1:].strip()
                # Preserve indentation for sub-items
                leading_spaces = len(line) - len(line.lstrip())
                if leading_spaces > 0:
                    bullet_text = "  " + bullet_text  # Mark as sub-item
                if current_slide:
                    current_slide["content"].append(bullet_text)
                else:
                    # Create a default content slide if none exists
                    current_slide = {
                        "type": "content",
                        "title": "Content",
                        "content": [bullet_text],
                    }
            # Handle other lines as content
            elif current_slide:
                # Could be a continuation or standalone text
                if current_slide["content"]:
                    # Append to last item or as new bullet
                    current_slide["content"].append(line_stripped)
                else:
                    current_slide["content"].append(line_stripped)

        # Add last slide
        if current_slide:
            slides.append(current_slide)

        # If no title slide was found but title is provided, add one
        if not has_title_slide and title:
            slides.insert(0, {"type": "title", "title": title, "content": []})

        # Ensure at least one slide exists
        if not slides:
            slides.append(
                {
                    "type": "title",
                    "title": title or "Presentation",
                    "content": [
                        content[:100] + "..." if len(content) > 100 else content
                    ],
                }
            )

        return slides

    def _get_template_colors(self, template: str) -> dict:
        """Get color scheme based on template style.

        Args:
            template: Template name

        Returns:
            Dict with title, text, and accent color RGB tuples
        """
        templates = {
            "default": {
                "title": (0, 112, 192),  # Blue
                "text": (51, 51, 51),  # Dark gray
                "accent": (0, 176, 240),  # Light blue
            },
            "professional": {
                "title": (31, 73, 125),  # Navy blue
                "text": (64, 64, 64),  # Gray
                "accent": (79, 129, 189),  # Steel blue
            },
            "minimal": {
                "title": (34, 34, 34),  # Near black
                "text": (68, 68, 68),  # Medium gray
                "accent": (102, 102, 102),  # Light gray
            },
        }
        return templates.get(template.lower(), templates["default"])

    def _escape_string(self, s: str) -> str:
        """Escape string for Python code generation.

        Args:
            s: Input string

        Returns:
            Escaped string safe for Python triple-quoted strings
        """
        # Replace backslashes first
        s = s.replace("\\", "\\\\")
        # Replace triple quotes
        s = s.replace('"""', '\\"\\"\\"')
        return s

    def _estimate_slide_count(self, content: str) -> int:
        """Estimate number of slides from content.

        Args:
            content: Content outline

        Returns:
            Estimated slide count
        """
        count = 0
        lines = content.lower().split("\n")
        for line in lines:
            if "title slide:" in line or line.strip().startswith("slide"):
                count += 1
        return max(count, 1)

    async def _execute_in_sandbox(self, code: str) -> Optional[bytes]:
        """Execute python-pptx code in E2B sandbox and download result.

        Args:
            code: Python code to execute

        Returns:
            Generated PPTX file bytes or None on failure
        """
        logger.info("[PPTXTool] Executing code in E2B sandbox...")

        try:
            from e2b_code_interpreter import Sandbox

            # Create sandbox
            e2b_api_key = os.getenv("E2B_API_KEY")
            if not e2b_api_key:
                logger.error("[PPTXTool] E2B_API_KEY not configured")
                raise ValueError("E2B_API_KEY environment variable is not set")

            loop = asyncio.get_running_loop()

            # Create sandbox in thread pool
            sandbox = await loop.run_in_executor(
                None,
                lambda: Sandbox(api_key=e2b_api_key, timeout=self.timeout),
            )

            try:
                # Install python-pptx if needed
                logger.debug("[PPTXTool] Installing python-pptx...")
                await loop.run_in_executor(
                    None,
                    lambda: sandbox.commands.run("pip install python-pptx pillow"),
                )

                # Execute the code
                logger.debug("[PPTXTool] Executing PPTX generation code...")
                execution = await loop.run_in_executor(
                    None,
                    lambda: sandbox.run_code(code),
                )

                # Check for errors
                if execution.error:
                    error_msg = f"{execution.error.name}: {execution.error.value}"
                    logger.warning(f"[PPTXTool] Execution error: {error_msg}")
                    raise RuntimeError(error_msg)

                # Check output for success message
                if execution.text and "SUCCESS:" in execution.text:
                    logger.info("[PPTXTool] Code execution successful")
                else:
                    logger.warning(f"[PPTXTool] Unexpected output: {execution.text}")

                # Download the generated file
                logger.debug(f"[PPTXTool] Downloading file from {DEFAULT_OUTPUT_PATH}")
                file_content = await loop.run_in_executor(
                    None,
                    lambda: sandbox.files.read(DEFAULT_OUTPUT_PATH),
                )

                if file_content:
                    # file_content is returned as bytes
                    if isinstance(file_content, str):
                        file_bytes = file_content.encode("latin-1")
                    else:
                        file_bytes = file_content
                    logger.info(f"[PPTXTool] File downloaded: {len(file_bytes)} bytes")
                    return file_bytes
                else:
                    logger.warning("[PPTXTool] No file content returned")
                    return None

            finally:
                # Clean up sandbox
                await loop.run_in_executor(None, sandbox.kill)
                logger.debug("[PPTXTool] Sandbox closed")

        except ImportError as e:
            logger.exception("[PPTXTool] E2B SDK not available")
            raise RuntimeError(
                "E2B SDK not installed. Please install e2b-code-interpreter."
            ) from e
        except Exception as e:
            logger.error(f"[PPTXTool] Sandbox execution failed: {e}", exc_info=True)
            raise

    async def _store_attachment(
        self,
        file_bytes: bytes,
        filename: str,
    ) -> dict:
        """Store generated file as attachment in database.

        Args:
            file_bytes: File binary data
            filename: Filename for the attachment

        Returns:
            Dict with success status, attachment_id, and download_url
        """
        logger.info(
            f"[PPTXTool] Storing attachment: filename={filename}, "
            f"size={len(file_bytes)} bytes, subtask_id={self.subtask_id}"
        )

        try:
            # Import context service and database utilities
            from app.db.session import SessionLocal
            from app.services.context import context_service

            # Create database session
            db = SessionLocal()

            try:
                # Use context service to upload attachment
                context, _truncation_info = context_service.upload_attachment(
                    db=db,
                    user_id=self.user_id,
                    filename=filename,
                    binary_data=file_bytes,
                    subtask_id=self.subtask_id,
                )

                attachment_id = context.id
                download_url = f"/api/v1/attachments/{attachment_id}/download"

                logger.info(
                    f"[PPTXTool] Attachment stored: id={attachment_id}, "
                    f"url={download_url}"
                )

                return {
                    "success": True,
                    "attachment_id": attachment_id,
                    "download_url": download_url,
                }

            finally:
                db.close()

        except ImportError as e:
            # Running in HTTP mode without backend access
            logger.warning(
                f"[PPTXTool] Backend services not available (HTTP mode): {e}"
            )
            # In HTTP mode, return the file info without storage
            # The caller should handle this appropriately
            return {
                "success": False,
                "error": "Attachment storage not available in HTTP mode",
            }
        except Exception as e:
            logger.error(
                f"[PPTXTool] Failed to store attachment: {e}",
                exc_info=True,
            )
            return {
                "success": False,
                "error": f"Failed to store attachment: {e!s}",
            }

    async def _emit_tool_status(
        self, status: str, message: str = "", result: Optional[dict] = None
    ) -> None:
        """Emit tool status update to frontend via WebSocket.

        Args:
            status: Status string ("running", "completed", "failed")
            message: Optional status message
            result: Optional result data for completed status
        """
        if not self.ws_emitter:
            return

        try:
            tool_output = {"message": message}
            if result:
                tool_output.update(result)

            await self.ws_emitter.emit_tool_call(
                task_id=self.task_id,
                tool_name=self.name,
                tool_input={},
                tool_output=tool_output,
                status=status,
            )
        except Exception as e:
            logger.warning(f"[PPTXTool] Failed to emit tool status: {e}")
